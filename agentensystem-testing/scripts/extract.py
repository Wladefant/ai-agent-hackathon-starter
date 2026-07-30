"""
Document Extractor - Converts documents to normalized markdown
NO LLM required - pure text extraction
Each input file → one markdown file in output/extracted/
Includes automatic domain detection
"""
import os
import sys
from pathlib import Path
from datetime import datetime
import json
import re
import inspect

# Ensure UTF-8 output on consoles that default to cp1252 (Windows PowerShell).
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

# Domain detection patterns
DOMAIN_PATTERNS = {
    "Core Banking/Payments": [
        r"\bSEPA\b", r"\bSWIFT\b", r"\bISO\s*20022\b", r"\bpacs\.\d{3}\b", 
        r"\bcamt\.\d{3}\b", r"\bIBAN\b", r"\bBIC\b", r"\binstant\s+payment\b",
        r"\bSCT\s*Inst\b", r"\bTIPS\b", r"\bRT1\b", r"\bEBA\s*Clearing\b",
        r"\bpain\.\d{3}\b", r"\bSEPA\s*Inst\b", r"\bpayment\s+initiation\b",
        r"\badmi\.\d{3}\b", r"\bcredit\s+transfer\b"
    ],
    "Insurance": [
        r"\bpolicy\b", r"\bclaim\b", r"\bpremium\b", r"\bunderwriting\b",
        r"\bcoverage\b", r"\binsured\b", r"\bbeneficiary\b", r"\bactuarial\b"
    ],
    "Lending": [
        r"\bloan\b", r"\bmortgage\b", r"\bcredit\s+score\b", r"\bdisbursement\b",
        r"\bamortization\b", r"\bcollateral\b", r"\binterest\s+rate\b"
    ],
    "Trade Finance": [
        r"\bletter\s+of\s+credit\b", r"\bLC\b", r"\bdocumentary\b",
        r"\btrade\s+finance\b", r"\bexport\b", r"\bimport\b", r"\bguarantee\b"
    ],
    "Wealth Management": [
        r"\bportfolio\b", r"\binvestment\b", r"\basset\s+management\b",
        r"\bcustody\b", r"\bsecurities\b", r"\bfund\b", r"\bNAV\b"
    ],
}

def detect_domain(text: str, filename: str = "") -> tuple[str, float]:
    """
    Detect domain from text content (and the file name, which often carries
    strong domain signals such as "Instant Payment").
    Returns (domain_name, confidence_score)
    """
    haystack = f"{filename}\n{text}"
    scores = {}

    for domain, patterns in DOMAIN_PATTERNS.items():
        matches = 0
        for pattern in patterns:
            if re.search(pattern, haystack, re.IGNORECASE):
                matches += 1
        if matches > 0:
            scores[domain] = matches / len(patterns)

    if scores:
        best_domain = max(scores, key=scores.get)
        return best_domain, scores[best_domain]

    return "Generic/IT", 0.0

def extract_pdf(file_path: str, output_dir: str) -> str:
    """Extract text and images from PDF files."""
    try:
        import fitz  # PyMuPDF
        from pathlib import Path
        
        doc = fitz.open(file_path)
        text_parts = []
        
        # Create a directory for images relative to the output markdown file
        pdf_path = Path(file_path)
        image_folder_name = f"{pdf_path.stem}_images"
        image_output_dir = Path(output_dir) / image_folder_name
        image_output_dir.mkdir(parents=True, exist_ok=True)
        
        for i, page in enumerate(doc):
            # Extract text
            page_text = page.get_text() or ""
            if page_text.strip():
                text_parts.append(f"## Page {i + 1}\n\n{page_text}")
            
            # Extract images
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                image_filename = f"page{i+1}_img{img_index}.{image_ext}"
                image_save_path = image_output_dir / image_filename
                
                with open(image_save_path, "wb") as img_file:
                    img_file.write(image_bytes)
                
                # Add image reference to markdown
                # Use relative path for portability
                relative_image_path = f"./{image_folder_name}/{image_filename}"
                text_parts.append(f"\n![Image from page {i+1}]({relative_image_path})\n")

        return "\n\n".join(text_parts)
    except ImportError:
        return "[ERROR: PyMuPDF (fitz) not installed. Run: pip install pymupdf]"
    except Exception as e:
        return f"[ERROR extracting PDF: {e}]"

def extract_docx(file_path: str) -> str:
    """Extract text from DOCX files."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                if para.style.name.startswith('Heading'):
                    level = para.style.name[-1] if para.style.name[-1].isdigit() else '2'
                    text_parts.append(f"{'#' * int(level)} {para.text}")
                else:
                    text_parts.append(para.text)
        for table in doc.tables:
            if table.rows:
                header_cells = [cell.text.strip() for cell in table.rows[0].cells]
                table_md = "\n| " + " | ".join(header_cells) + " |\n"
                table_md += "| " + " | ".join(["---"] * len(header_cells)) + " |\n"
                for row in table.rows[1:]:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_md += f"| {row_text} |\n"
                text_parts.append(table_md)
        return "\n\n".join(text_parts)
    except ImportError:
        return "[ERROR: python-docx not installed. Run: pip install python-docx]"
    except Exception as e:
        return f"[ERROR extracting DOCX: {e}]"

def extract_pptx(file_path: str) -> str:
    """Extract text from PPTX files."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"## Slide {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            text_parts.append("\n\n".join(slide_text))
        return "\n\n---\n\n".join(text_parts)
    except ImportError:
        return "[ERROR: python-pptx not installed. Run: pip install python-pptx]"
    except Exception as e:
        return f"[ERROR extracting PPTX: {e}]"

def extract_xlsx(file_path: str) -> str:
    """Extract text from XLSX files."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = [f"## Sheet: {sheet_name}\n"]
            rows = list(ws.iter_rows(values_only=True))
            if rows:
                header = rows[0] if rows else []
                header_clean = [str(h) if h else "" for h in header]
                sheet_text.append("| " + " | ".join(header_clean) + " |")
                sheet_text.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in rows[1:]:
                    row_clean = [str(c) if c else "" for c in row]
                    sheet_text.append("| " + " | ".join(row_clean) + " |")
            text_parts.append("\n".join(sheet_text))
        return "\n\n".join(text_parts)
    except ImportError:
        return "[ERROR: openpyxl not installed. Run: pip install openpyxl]"
    except Exception as e:
        return f"[ERROR extracting XLSX: {e}]"

def extract_txt(file_path: str) -> str:
    """Extract text from TXT files."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"[ERROR extracting TXT: {e}]"

def extract_md(file_path: str) -> str:
    """Copy markdown files as-is."""
    return extract_txt(file_path)


# ---------------------------------------------------------------------------
# Visio (VSDX) structured flow extraction
# ---------------------------------------------------------------------------
# Unlike a flat text dump, this reconstructs the diagram as a directed graph:
# shapes become nodes (with absolute geometry), and the <Connects> section is
# parsed into directed source -> target edges. Connector labels (Yes / No /
# error codes) become edge labels, and container shapes become swimlane groups.
# The result is emitted as a Mermaid flowchart + readable flow steps so the LLM
# reads the diagram AS A FLOW, not just a bag of labels.

_VISIO_NS = "http://schemas.microsoft.com/office/visio/2012/main"


def _vx_cell(shape, name):
    """Return a numeric Cell value (PinX, Width, ...) from a Visio shape."""
    for c in shape.findall(f"{{{_VISIO_NS}}}Cell"):
        if c.get("N") == name:
            try:
                return float(c.get("V"))
            except (TypeError, ValueError):
                return None
    return None


def _vx_text(shape):
    """Collapse a shape's <Text> run into a single clean string."""
    t = shape.find(f"{{{_VISIO_NS}}}Text")
    if t is None:
        return ""
    return re.sub(r"\s+", " ", "".join(t.itertext())).strip()


def _vx_is_container(shape):
    """True if the shape is a Visio container/swimlane (msvSDContainerLocked)."""
    user = shape.find(f'{{{_VISIO_NS}}}Section[@N="User"]')
    if user is None:
        return False
    for row in user.findall(f"{{{_VISIO_NS}}}Row"):
        if row.get("N") == "msvSDContainerLocked":
            return True
    return False


def _vx_is_connector(name):
    return name.startswith("Dynamic connector") or name.startswith("Connector")


def _vx_clean_label(text):
    """P4 - drop shape-ID noise ('02'), pure numbers/punctuation, and blanks."""
    t = (text or "").strip()
    if not t:
        return ""
    if t == "02" or re.fullmatch(r"[\d\W_]+", t):
        return ""
    return t


def _vx_walk(shape, parent_ox, parent_oy, depth, nodes):
    """Recursively record every shape with its ABSOLUTE centre coordinates.

    Visio child coordinates are relative to the parent group's local origin,
    so absolute positions are accumulated down the shape tree. Rotation is
    ignored (flowcharts are axis-aligned).
    """
    pinx = _vx_cell(shape, "PinX") or 0.0
    piny = _vx_cell(shape, "PinY") or 0.0
    locx = _vx_cell(shape, "LocPinX") or 0.0
    locy = _vx_cell(shape, "LocPinY") or 0.0
    width = _vx_cell(shape, "Width") or 0.0
    height = _vx_cell(shape, "Height") or 0.0
    cx = parent_ox + pinx
    cy = parent_oy + piny
    origin_x = cx - locx
    origin_y = cy - locy
    sid = shape.get("ID")
    name = shape.get("NameU", "") or ""
    if sid is not None:
        nodes[sid] = {
            "name": name,
            "kind": name.split(".")[0] if name else "Shape",
            "text": _vx_text(shape),
            "cx": cx, "cy": cy, "w": width, "h": height,
            "container": _vx_is_container(shape),
            "connector": _vx_is_connector(name),
        }
    sub = shape.find(f"{{{_VISIO_NS}}}Shapes")
    if sub is not None:
        for child in sub.findall(f"{{{_VISIO_NS}}}Shape"):
            _vx_walk(child, origin_x, origin_y, depth + 1, nodes)


def _vx_resolve_label(sid, nodes, max_dist=1.0):
    """Resolved label for a node: its own text, else the nearest labelled
    process shape WITHIN `max_dist` (Visio inches). Decision diamonds carry no
    text of their own; the nearest-neighbour heuristic only applies when a
    labelled shape is genuinely close, so distant shapes are not mislabelled."""
    n = nodes.get(sid)
    if not n:
        return ""
    own = _vx_clean_label(n["text"])
    if own:
        return own
    best, best_d = "", float("inf")
    for o in nodes.values():
        if o["container"] or o["connector"]:
            continue
        lt = _vx_clean_label(o["text"])
        if not lt:
            continue
        d = (o["cx"] - n["cx"]) ** 2 + (o["cy"] - n["cy"]) ** 2
        if d < best_d:
            best_d, best = d, lt
    if best and best_d ** 0.5 <= max_dist:
        return best
    return ""


def _vx_parse_page(xml_bytes):
    """Parse one Visio page into (nodes, edges, containers)."""
    import xml.etree.ElementTree as ET

    root = ET.fromstring(xml_bytes)
    nodes = {}
    top = root.find(f"{{{_VISIO_NS}}}Shapes")
    if top is not None:
        for sh in top.findall(f"{{{_VISIO_NS}}}Shape"):
            _vx_walk(sh, 0.0, 0.0, 0, nodes)

    # Edges from the <Connects> section: group each connector's BeginX (source)
    # and EndX (target) endpoints by the connector sheet ID.
    text = xml_bytes.decode("utf-8", errors="replace")
    endpoints = {}
    for m in re.finditer(
        r"<Connect\s+FromSheet='(\d+)'\s+FromCell='(\w+)'[^>]*?ToSheet='(\d+)'", text
    ):
        conn_id, cell_name, to_sheet = m.group(1), m.group(2), m.group(3)
        ep = endpoints.setdefault(conn_id, {})
        if cell_name == "BeginX":
            ep["src"] = to_sheet
        elif cell_name == "EndX":
            ep["dst"] = to_sheet

    edges = []
    for conn_id, ep in endpoints.items():
        src, dst = ep.get("src"), ep.get("dst")
        if not src or not dst or src == dst:
            continue
        edges.append({
            "src": src,
            "dst": dst,
            "label": _vx_clean_label(nodes.get(conn_id, {}).get("text", "")),
        })

    # Containers (swimlanes): name by the top-most labelled member inside the box.
    containers = []
    for c in nodes.values():
        if not c["container"]:
            continue
        name, top_cy = "", float("-inf")
        members = []
        for o in nodes.values():
            if o is c or o["container"] or o["connector"]:
                continue
            if abs(o["cx"] - c["cx"]) <= c["w"] / 2 and abs(o["cy"] - c["cy"]) <= c["h"] / 2:
                lt = _vx_clean_label(o["text"])
                if lt:
                    members.append(lt)
                    if o["cy"] > top_cy:
                        top_cy, name = o["cy"], lt
        if members:
            containers.append({"name": name or "(unnamed region)", "members": members})

    return nodes, edges, containers


def _vx_mermaid_id(sid):
    return f"n{sid}"


def _vx_mermaid_label(text):
    """Make a label safe for a Mermaid node/edge. Mermaid treats <, >, {, },
    |, " specially, so neutralise them (<> would otherwise be parsed as HTML)."""
    t = (text or "").replace("\n", " ").strip()
    t = (t.replace("<", "\u2039").replace(">", "\u203a")
          .replace("{", "(").replace("}", ")")
          .replace("|", "/").replace('"', "'"))
    t = re.sub(r"\s+", " ", t)
    return t[:80]


def extract_vsdx(file_path: str, output_dir: str) -> str:
    """Extract a Visio VSDX as a structured flow (nodes + directed edges).

    Emits, per page: a Mermaid flowchart, a readable list of flow steps, the
    swimlane/container grouping, and a de-duplicated label inventory. Falls back
    to a flat label dump only if no connectors can be parsed.
    """
    try:
        import zipfile

        result = [
            "> **Note:** Visio diagram. The flow below was reconstructed from "
            "shapes and their connectors (source -> target), not just loose text.",
            "",
        ]

        with zipfile.ZipFile(file_path, "r") as z:
            page_files = sorted(
                n for n in z.namelist()
                if re.match(r"visio/pages/page\d+\.xml$", n)
            )
            if not page_files:
                return "[No Visio pages found in VSDX]"

            any_edges = False
            for idx, page_name in enumerate(page_files, start=1):
                xml_bytes = z.read(page_name)
                nodes, edges, containers = _vx_parse_page(xml_bytes)
                labels = {
                    sid: _vx_resolve_label(sid, nodes)
                    for sid in nodes
                    if not nodes[sid]["connector"] and not nodes[sid]["container"]
                }

                result.append(f"## Page {idx}")
                result.append("")

                if edges:
                    any_edges = True
                    # --- Mermaid flowchart ---
                    referenced = set()
                    for e in edges:
                        referenced.add(e["src"])
                        referenced.add(e["dst"])
                    result.append("### Flow Diagram")
                    result.append("")
                    result.append("```mermaid")
                    result.append("flowchart TD")
                    for sid in referenced:
                        n = nodes.get(sid)
                        lbl = _vx_mermaid_label(labels.get(sid) or (n["kind"] if n else "node"))
                        if n and n["kind"].lower().startswith("decision"):
                            result.append(f'    {_vx_mermaid_id(sid)}{{"{lbl}"}}')
                        else:
                            result.append(f'    {_vx_mermaid_id(sid)}["{lbl}"]')
                    for e in edges:
                        a, b = _vx_mermaid_id(e["src"]), _vx_mermaid_id(e["dst"])
                        if e["label"]:
                            result.append(f'    {a} -->|"{_vx_mermaid_label(e["label"])}"| {b}')
                        else:
                            result.append(f"    {a} --> {b}")
                    result.append("```")
                    result.append("")

                    # --- Readable flow steps ---
                    result.append("### Flow Steps")
                    result.append("")

                    def _step_label(node_id):
                        lbl = labels.get(node_id)
                        if lbl:
                            return lbl
                        n = nodes.get(node_id)
                        return f"({n['kind']})" if n else "(unlabelled)"

                    for e in edges:
                        s = _step_label(e["src"])
                        d = _step_label(e["dst"])
                        if e["label"]:
                            result.append(f"- {s} --[{e['label']}]--> {d}")
                        else:
                            result.append(f"- {s} --> {d}")
                    result.append("")
                else:
                    # Fallback: flat label dump so we never regress to nothing.
                    flat = sorted({
                        _vx_clean_label(n["text"]) for n in nodes.values()
                    } - {""})
                    if flat:
                        result.append("### Diagram Elements")
                        result.append("")
                        result.extend(f"- {t}" for t in flat)
                        result.append("")

                # --- Swimlanes / containers ---
                if containers:
                    result.append("### Swimlanes / Containers")
                    result.append("")
                    for c in containers:
                        uniq = list(dict.fromkeys(c["members"]))
                        result.append(f"- **{c['name']}**: " + "; ".join(uniq[:20]))
                    result.append("")

                # --- Label inventory (also feeds domain detection) ---
                inv = sorted({v for v in labels.values() if v})
                if inv:
                    result.append("### Process / Activity Labels")
                    result.append("")
                    result.extend(f"- {t}" for t in inv)
                    result.append("")

            if not any_edges:
                result.append(
                    "> No connectors were found; only shape labels could be extracted."
                )

        return "\n".join(result)

    except Exception as e:
        return f"[ERROR extracting VSDX: {e}]"


EXTRACTORS = {
    '.pdf':   extract_pdf,
    '.docx':  extract_docx,
    '.pptx':  extract_pptx,
    '.xlsx':  extract_xlsx,
    '.txt':   extract_txt,
    '.md':    extract_md,
    '.vsdx':  extract_vsdx,
}

def extract_document(file_path: str, output_dir: str) -> str:
    """Extract text from a document based on its extension."""
    ext = Path(file_path).suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return f"[Unsupported file format: {ext}]"
    
    # Check if the extractor needs the output directory
    sig = inspect.signature(extractor)
    if 'output_dir' in sig.parameters:
        return extractor(file_path, output_dir)
    else:
        # For other extractors that don't need output_dir
        return extractor(file_path)

def process_inputs(input_dir: str = "inputs", output_dir: str = "extracted"):
    """Process all documents - one input file = one markdown output."""
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    
    if not input_path.exists() or not input_path.is_dir():
        print(f"ERROR: Input directory '{input_dir}' not found or is not a directory.")
        sys.exit(1)
        
    all_files = list(input_path.rglob("*.*"))
    supported_files = [f for f in all_files if f.suffix.lower() in EXTRACTORS]
    
    if not supported_files:
        print(f"ERROR: No supported files found in '{input_dir}'.")
        sys.exit(1)
        
    manifest = {
        "files": []
    }

    for root, _, files in os.walk(input_path):
        for file in files:
            file_path = Path(root) / file
            if file_path.suffix.lower() not in EXTRACTORS:
                continue

            relative_path = file_path.relative_to(input_path)
            
            # Create a corresponding output directory structure
            md_output_dir = output_path / relative_path.parent
            md_output_dir.mkdir(parents=True, exist_ok=True)
            
            md_filename = file_path.stem + ".md"
            md_output_path = md_output_dir / md_filename

            print(f"Processing: {relative_path}")
            
            # Pass the specific output directory for this file to extract_document
            content = extract_document(str(file_path), str(md_output_dir))
            
            domain, confidence = detect_domain(content, file_path.name)
            
            # Create frontmatter
            frontmatter = (
                f"---\n"
                f"source: {file_path.name}\n"
                f"source_path: {str(relative_path)}\n"
                f"extracted: {datetime.now().isoformat()}\n"
                f"type: {file_path.suffix}\n"
                f"domain: {domain}\n"
                f"domain_confidence: {confidence:.2f}\n"
                f"---\n\n"
            )
            
            full_content = frontmatter + f"# {file_path.stem}\n\n" + content
            
            with open(md_output_path, 'w', encoding='utf-8') as f:
                f.write(full_content)
            
            print(f"  → {md_output_path} (Domain: {domain})")
            
            manifest["files"].append({
                "source": str(relative_path),
                "output": str(md_output_path.relative_to(output_path)),
                "domain": domain,
                "confidence": confidence
            })

    # Save manifest
    manifest_path = output_path / (input_path.name + "_extraction_manifest.json")
    
    # Consolidate domain detection
    overall_domain, overall_confidence = consolidate_domains(manifest)
    manifest['detected_domain'] = overall_domain
    manifest['domain_confidence'] = overall_confidence

    with open(manifest_path, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, indent=4)

    print(f"\n✅ Processed {len(manifest['files'])} documents.")
    print(f"📄 Manifest: {manifest_path}")
    print(f"🎯 Detected Domain: {overall_domain} (confidence: {overall_confidence:.0%})")
    print("\n✅ Extraction complete!")

def consolidate_domains(manifest: dict) -> tuple[str, float]:
    """Determine the most likely domain from all processed files."""
    domain_scores = {}
    total_confidence = 0
    
    if not manifest["files"]:
        return "Generic/IT", 0.0

    for file in manifest["files"]:
        domain = file["domain"]
        confidence = file["confidence"]
        
        if domain != "Generic/IT":
            domain_scores[domain] = domain_scores.get(domain, 0) + confidence
            total_confidence += confidence
            
    if not domain_scores:
        return "Generic/IT", 0.0
        
    # Normalize scores
    best_domain = max(domain_scores, key=domain_scores.get)
    
    # The overall confidence could be the average or the max score's proportion
    overall_confidence = domain_scores[best_domain] / total_confidence if total_confidence > 0 else 0
    
    return best_domain, overall_confidence

def main():
    """Main entry point for the script."""
    import argparse
    parser = argparse.ArgumentParser(description="Extract text and metadata from documents.")
    parser.add_argument("--input", default="inputs", help="Input directory containing documents.")
    parser.add_argument("--output", default="output/extracted", help="Output directory for markdown files.")
    args = parser.parse_args()
    
    process_inputs(args.input, args.output)

if __name__ == "__main__":
    main()
